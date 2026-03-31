"""
PDF → 完全再現 PPTX 変換スクリプト
使い方: python pdf_to_pptx.py

方針:
  1. PDFページを高解像度画像として背景に使用（イラスト保持）
  2. numpy でピンクテキストボックスを自動検出 → 白マスク
  3. 検出位置に編集可能なピンク角丸ボックス＋テキストを配置
"""

import fitz  # PyMuPDF
from PIL import Image
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import io

PDF_PATH   = r"C:\Users\Owner\Downloads\Gentle_AI_Support.pdf"
OUTPUT_PATH = r"C:\Users\Owner\Downloads\Gentle_AI_Support_editable.pptx"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
RENDER_SCALE = 3.0   # 高解像度レンダリング倍率

# カラー
PINK_ACCENT   = RGBColor(0xe0, 0x7b, 0x8a)
PINK_LIGHT_BG = RGBColor(0xfd, 0xf0, 0xf2)
TEXT_DARK     = RGBColor(0x33, 0x33, 0x33)
WHITE         = RGBColor(0xff, 0xff, 0xff)

# 各スライドのテキスト
SLIDES = [
    {
        "title": "AIを、心強い味方にするために",
        "body": [
            "思春期の子育て、毎日本当にお疲れ様です",
            "今日は、お母さん自身の心をフワッと軽くするお話です",
        ]
    },
    {
        "title": "今日のテーマは\n「お母さんの心のケア」",
        "body": [
            "毎日、子どものことばかり優先していませんか？",
            "まずは、あなた自身の素直な気持ちに耳を傾けましょう",
        ]
    },
    {
        "title": "自分のモヤモヤを\n「言葉にする」練習",
        "body": [
            "心の中のモヤモヤは、言葉にすると不思議と整理されます",
            "無意識に浮かぶ感情に「気づく」ことが第一歩です",
        ]
    },
    {
        "title": "こんなふうに自分を\n責めていませんか？",
        "body": [
            "「また怒っちゃった… 私のせいかな」",
            "これは、とっさに浮かぶネガティブな思い込み（心のクセ）です",
        ]
    },
    {
        "title": "そこで、AIを「練習相手」にしてみる",
        "body": [
            "家族には言いにくいことも、AIなら気を遣いません",
            "いつでも、どんなことでも、受け止めてくれます",
        ]
    },
    {
        "title": "お話し相手になる\n「傾聴AI」とは？",
        "body": [
            "あなたのお話を「ただ聞いてくれる」無料のツールです",
            "アドバイスや否定はせず、優しく寄り添うように設定されています",
        ]
    },
    {
        "title": "使い方はとってもシンプル",
        "body": [
            "LINEでお友達にメッセージを送るような感覚でOK",
            "チャット画面に、今の気持ちをそのまま打ち込むだけです",
        ]
    },
    {
        "title": "大切なルール①\n主役は「あなた」です",
        "body": [
            "AIの返事に、無理に合わせる必要はありません",
            "「心のハンドル」は自分で握り、自分のペースで使いましょう",
        ]
    },
    {
        "title": "大切なルール②\n頼る先はひとつじゃない",
        "body": [
            "AIは、あくまで心を軽くするサポート役のひとつです",
            "家族、友人、趣味の時間など、色々な「ホッとする居場所」を持ちましょう",
        ]
    },
    {
        "title": "大切なルール③ 辛いときは専門家へ",
        "body": [
            "心がどうしても苦しいときは、決して無理をしないでください",
            "AIではなく、人間のお医者さんやカウンセラーに頼りましょう",
        ]
    },
    {
        "title": "個人情報には気をつけて",
        "body": [
            "AIは安全ですが、念のため秘密はしっかり守りましょう",
            "本名や住所など、個人が特定できる情報は入力しないでくださいね",
        ]
    },
    {
        "title": "今日からできる、小さな第一歩",
        "body": [
            "まずは今日感じたことを、AIに一言だけ送ってみませんか？",
            "「あー、今日は疲れた！」の一言だけでも大歓迎です",
        ]
    },
    {
        "title": "AIは、あなたの「秘密のノート」",
        "body": [
            "誰にも見られない、あなただけの心のための練習ノートです",
            "毎日がんばるお母さんの心が、少しでもフワッと軽くなりますように",
        ]
    },
]


# ─────────────────────────────────────────────
# ピンクボックス検出＆マスク
# ─────────────────────────────────────────────

def detect_and_mask_pink_boxes(pil_img):
    """
    PIL画像から淡いピンクのテキストボックス領域を検出し、
    白でマスクした画像と検出したボックス一覧（fractions）を返す。
    """
    arr = np.array(pil_img, dtype=np.int32)
    R, G, B = arr[:, :, 0], arr[:, :, 1], arr[:, :, 2]
    H, W = R.shape

    # 淡いピンク判定（テキストボックスの背景色 #fdf0f2 ≈ R:253 G:240 B:242）
    is_pink = (
        (R > 240) &
        (G > 222) &
        (B > 225) &
        ((R - G) > 4) &          # わずかに赤みがかっている
        ((R - B) > 2) &
        ~((R > 252) & (G > 252) & (B > 252))  # 白は除外
    )

    # 小さなノイズを除去（行・列方向に走査して大きな矩形領域のみ残す）
    MIN_BOX_AREA_FRAC = 0.008  # スライド面積の 0.8% 以上

    boxes_px = _find_large_rects(is_pink, H, W, MIN_BOX_AREA_FRAC)

    # マスク処理（ボックス領域を白で塗りつぶす）
    masked = pil_img.copy()
    from PIL import ImageDraw
    draw = ImageDraw.Draw(masked)
    for (r0, c0, r1, c1) in boxes_px:
        pad = 6
        draw.rectangle([c0 - pad, r0 - pad, c1 + pad, r1 + pad], fill=(255, 255, 255))

    # 座標をスライド比率に変換
    boxes_frac = [(r0/H, c0/W, r1/H, c1/W) for (r0, c0, r1, c1) in boxes_px]
    return masked, boxes_frac


def _find_large_rects(mask, H, W, min_area_frac):
    """
    boolean mask から大きな矩形領域を探す（シンプルな行投影法）。
    Returns list of (row0, col0, row1, col1) in pixel coords.
    """
    min_area = H * W * min_area_frac

    # 行ごとの pink ピクセル数
    row_sum = mask.sum(axis=1)   # shape: (H,)
    col_sum = mask.sum(axis=0)   # shape: (W,)

    # 有効行・列のしきい値（全幅/高さの 8% 以上が pink）
    row_thresh = W * 0.06
    col_thresh = H * 0.06

    active_rows = np.where(row_sum > row_thresh)[0]
    active_cols = np.where(col_sum > col_thresh)[0]

    if len(active_rows) == 0 or len(active_cols) == 0:
        return []

    # 連続区間にグループ化
    row_groups = _group_consecutive(active_rows, gap=8)
    col_groups = _group_consecutive(active_cols, gap=8)

    boxes = []
    for rg in row_groups:
        r0, r1 = rg[0], rg[-1]
        for cg in col_groups:
            c0, c1 = cg[0], cg[-1]
            # そのサブ領域でピンクが多いか確認
            sub = mask[r0:r1+1, c0:c1+1]
            pink_ratio = sub.sum() / (sub.size + 1e-9)
            area = (r1 - r0 + 1) * (c1 - c0 + 1)
            if pink_ratio > 0.35 and area > min_area:
                boxes.append((r0, c0, r1, c1))

    return boxes


def _group_consecutive(arr, gap=5):
    """連続した整数配列をグループに分割"""
    if len(arr) == 0:
        return []
    groups = []
    cur = [arr[0]]
    for v in arr[1:]:
        if v - cur[-1] <= gap:
            cur.append(v)
        else:
            groups.append(cur)
            cur = [v]
    groups.append(cur)
    return groups


# ─────────────────────────────────────────────
# PPTX ユーティリティ
# ─────────────────────────────────────────────

def add_rounded_rect_shape(slide, left_in, top_in, width_in, height_in,
                            fill_rgb, border_rgb, border_pt=1.5):
    """角丸四角形を追加"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle（後でroundRectに変換）
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    sp = shape.element
    spPr = sp.find(qn("p:spPr"))
    prstGeom = spPr.find(qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "roundRect")
        avLst = prstGeom.find(qn("a:avLst"))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn("a:avLst"))
        for gd in avLst.findall(qn("a:gd")):
            avLst.remove(gd)
        gd_el = etree.SubElement(avLst, qn("a:gd"))
        gd_el.set("name", "adj")
        gd_el.set("fmla", "val 8000")  # 8% 角丸

    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = border_rgb
    shape.line.width = Pt(border_pt)
    return shape


def set_shape_text(shape, lines, font_pt, color_rgb, bold=False, align=PP_ALIGN.LEFT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_top    = Pt(10)
    tf.margin_bottom = Pt(10)
    tf.margin_left   = Pt(18)
    tf.margin_right  = Pt(18)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(font_pt)
        run.font.color.rgb = color_rgb
        run.font.bold = bold


def add_title_box(slide, text, left_in, top_in, width_in, height_in, font_pt=38):
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(font_pt)
        run.font.bold = True
        run.font.color.rgb = TEXT_DARK
    tb.fill.background()
    tb.line.fill.background()
    return tb


# ─────────────────────────────────────────────
# スライドレイアウト定義
# ─────────────────────────────────────────────
# title_area: (left%, top%, width%, height%)
# body_areas: [(left%, top%, width%, height%), ...]  ← 本文ボックスごと
# ※ % は 0-1 の比率（スライドサイズに対する割合）

LAYOUTS = [
    # Slide 1: 女性（右） ← タイトル・ボックス左
    {
        "title_area": (0.02, 0.03, 0.53, 0.30),
        "body_areas": [
            (0.02, 0.50, 0.53, 0.19),
            (0.02, 0.72, 0.53, 0.19),
        ],
        "title_font": 36,
    },
    # Slide 2: ハート（中央）← タイトル上・ボックス下左右
    {
        "title_area": (0.08, 0.02, 0.84, 0.20),
        "body_areas": [
            (0.02, 0.64, 0.46, 0.30),
            (0.52, 0.64, 0.46, 0.30),
        ],
        "title_font": 36,
    },
    # Slide 3: 毛糸玉（左） ← タイトル左・ボックス右
    {
        "title_area": (0.02, 0.02, 0.44, 0.26),
        "body_areas": [
            (0.47, 0.28, 0.50, 0.26),
            (0.47, 0.57, 0.50, 0.26),
        ],
        "title_font": 34,
    },
    # Slide 4: 女性（左） ← タイトル上・ボックス右
    {
        "title_area": (0.04, 0.02, 0.92, 0.18),
        "body_areas": [
            (0.41, 0.27, 0.55, 0.27),
            (0.41, 0.57, 0.55, 0.29),
        ],
        "title_font": 36,
    },
    # Slide 5: お日様と人（中央）← タイトル上・ボックス下左右
    {
        "title_area": (0.02, 0.02, 0.96, 0.18),
        "body_areas": [
            (0.02, 0.65, 0.43, 0.30),
            (0.55, 0.65, 0.43, 0.30),
        ],
        "title_font": 36,
    },
    # Slide 6: スマホ（中央）← タイトル左・ボックス右
    {
        "title_area": (0.02, 0.02, 0.46, 0.30),
        "body_areas": [
            (0.48, 0.25, 0.49, 0.28),
            (0.48, 0.57, 0.49, 0.28),
        ],
        "title_font": 34,
    },
    # Slide 7: スマホ・手（左） ← タイトル上・ボックス右
    {
        "title_area": (0.18, 0.02, 0.64, 0.18),
        "body_areas": [
            (0.50, 0.19, 0.47, 0.31),
            (0.50, 0.54, 0.47, 0.31),
        ],
        "title_font": 36,
    },
    # Slide 8: ロープリング（右）← タイトル左・ボックス左下
    {
        "title_area": (0.02, 0.02, 0.54, 0.32),
        "body_areas": [
            (0.02, 0.57, 0.53, 0.17),
            (0.02, 0.76, 0.53, 0.17),
        ],
        "title_font": 34,
    },
    # Slide 9: 桜の木（右）← タイトル左・ボックス左下
    {
        "title_area": (0.02, 0.02, 0.54, 0.32),
        "body_areas": [
            (0.02, 0.45, 0.46, 0.19),
            (0.02, 0.66, 0.46, 0.29),
        ],
        "title_font": 34,
    },
    # Slide 10: 手と十字（中央）← タイトル上・ワイドボックス下
    {
        "title_area": (0.08, 0.02, 0.84, 0.16),
        "body_areas": [
            (0.02, 0.72, 0.96, 0.22),
        ],
        "title_font": 36,
        "wide_body": True,
    },
    # Slide 11: 鍵（右上）← タイトル左・ボックス中央
    {
        "title_area": (0.02, 0.02, 0.73, 0.18),
        "body_areas": [
            (0.07, 0.28, 0.81, 0.24),
            (0.07, 0.55, 0.81, 0.24),
        ],
        "title_font": 36,
    },
    # Slide 12: 足跡（中央）← タイトル上・テキスト下
    {
        "title_area": (0.08, 0.02, 0.84, 0.16),
        "body_areas": [
            (0.02, 0.63, 0.96, 0.17),
            (0.06, 0.81, 0.88, 0.14),
        ],
        "title_font": 36,
    },
    # Slide 13: ノート（中央）← タイトル上・テキスト下
    {
        "title_area": (0.04, 0.02, 0.92, 0.22),
        "body_areas": [
            (0.07, 0.71, 0.86, 0.13),
            (0.07, 0.83, 0.86, 0.13),
        ],
        "title_font": 36,
    },
]


# ─────────────────────────────────────────────
# メイン処理
# ─────────────────────────────────────────────

def build_pptx():
    doc = fitz.open(PDF_PATH)
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    W_IN = 13.33  # スライド幅 (inches)
    H_IN = 7.5    # スライド高さ (inches)

    for i, (slide_data, layout) in enumerate(zip(SLIDES, LAYOUTS)):
        slide = prs.slides.add_slide(blank_layout)

        # ── 1. PDFページを高解像度でレンダリング ──────────────
        page = doc[i]
        mat  = fitz.Matrix(RENDER_SCALE, RENDER_SCALE)
        pix  = page.get_pixmap(matrix=mat, alpha=False)
        img  = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")

        # ── 2. ピンクボックスをマスク（自動検出＋白塗り）─────
        masked_img, detected_boxes = detect_and_mask_pink_boxes(img)

        if detected_boxes:
            print(f"  スライド{i+1}: {len(detected_boxes)}個のボックス検出")
        else:
            print(f"  スライド{i+1}: ボックス未検出（レイアウト定義を使用）")

        # ── 3. マスク済み画像をスライド背景に貼る ────────────
        buf = io.BytesIO()
        masked_img.save(buf, format="PNG")
        buf.seek(0)
        slide.shapes.add_picture(buf, 0, 0, SLIDE_W, SLIDE_H)

        # ── 4. タイトルテキストボックス（透明背景） ───────────
        la, lt, lw, lh = layout["title_area"]
        font_pt = layout.get("title_font", 36)
        add_title_box(
            slide,
            slide_data["title"],
            la * W_IN, lt * H_IN,
            lw * W_IN, lh * H_IN,
            font_pt=font_pt
        )

        # ── 5. 本文：ピンク角丸ボックス ───────────────────────
        body_lines = slide_data["body"]
        body_areas = layout["body_areas"]
        wide_body  = layout.get("wide_body", False)

        if wide_body and len(body_lines) >= 2:
            # ワイドボックス: 2行を1つのボックスにまとめる
            ba = body_areas[0]
            shape = add_rounded_rect_shape(
                slide,
                ba[0] * W_IN, ba[1] * H_IN,
                ba[2] * W_IN, ba[3] * H_IN,
                PINK_LIGHT_BG, PINK_ACCENT
            )
            set_shape_text(shape, body_lines, 22, TEXT_DARK)
        else:
            for j, ba in enumerate(body_areas):
                if j >= len(body_lines):
                    break
                shape = add_rounded_rect_shape(
                    slide,
                    ba[0] * W_IN, ba[1] * H_IN,
                    ba[2] * W_IN, ba[3] * H_IN,
                    PINK_LIGHT_BG, PINK_ACCENT
                )
                set_shape_text(shape, [body_lines[j]], 22, TEXT_DARK)

        print(f"  スライド {i+1}/{len(SLIDES)} 完了")

    doc.close()
    prs.save(OUTPUT_PATH)
    print(f"\n保存完了: {OUTPUT_PATH}")


if __name__ == "__main__":
    print("変換中...\n")
    build_pptx()
